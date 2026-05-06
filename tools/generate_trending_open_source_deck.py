from __future__ import annotations

import base64
import json
import subprocess
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.presentation import Presentation as PresentationType
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_DIR = Path("artifacts")
PPTX_PATH = OUTPUT_DIR / "trending_open_source_daily.pptx"
VIDEO_PATH = OUTPUT_DIR / "trending_open_source_daily.mp4"
MANIFEST_PATH = OUTPUT_DIR / "trending_open_source_daily_manifest.json"
TIMELINE_MD_PATH = OUTPUT_DIR / "trending_open_source_daily_tts_timeline.md"

BG_TOP = RGBColor(10, 17, 40)
BG_BOTTOM = RGBColor(245, 247, 250)
NAVY = RGBColor(17, 24, 39)
BLUE = RGBColor(42, 91, 215)
CYAN = RGBColor(15, 180, 209)
GOLD = RGBColor(246, 185, 59)
SOFT = RGBColor(236, 241, 248)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(95, 106, 129)
WHITE = RGBColor(255, 255, 255)

ICON_FONT = "Segoe UI Emoji"

CARD_META = {
    "项目定位": ("🧭", "项目定位"),
    "核心卖点": ("🚀", "核心特点"),
    "场景": ("🎯", "适用场景"),
    "热度原因": ("🔥", "关注原因"),
}

SLIDES = [
    {
        "rank": 1,
        "repo": "mattpocock/skills",
        "language": "Shell",
        "stars_today": 350,
        "stars_total": 5200,
        "headline": "个人技能库变成了 AI 助手的能力市场",
        "summary": "把规划、PRD、重构和界面设计拆成可安装技能，正在推动 agent 工作流模块化。",
        "bullets": [
            "核心卖点: 可组合的 agent skills 目录",
            "场景: 产品规划、开发执行、重构拆解",
            "热度原因: 低门槛复用 Claude 风格能力包",
        ],
        "narration": "Skills 把高频开发动作封装成可安装技能，适合做自动化流程里的能力积木。",
        "duration": 10,
    },
    {
        "rank": 2,
        "repo": "jarrodwatts/claude-hud",
        "language": "JavaScript",
        "stars_today": 305,
        "stars_total": 7971,
        "headline": "给编码助手加上一块实时状态仪表盘",
        "summary": "直接显示上下文占用、工具执行、子代理和 todo 进度，解决 agent 可观测性问题。",
        "bullets": [
            "核心卖点: 状态栏可视化 agent 内部流程",
            "场景: 调试长任务、观察上下文消耗",
            "热度原因: 让黑箱式 CLI 助手更可控",
        ],
        "narration": "Claude HUD 的价值在可观测性，它把上下文、工具和任务进度实时暴露出来。",
        "duration": 10,
    },
    {
        "rank": 3,
        "repo": "THU-MAIC/OpenMAIC",
        "language": "TypeScript",
        "stars_today": 288,
        "stars_total": 8355,
        "headline": "多智能体课堂把学习体验做成沉浸式产品",
        "summary": "通过多角色协作与一键部署，构建可以讲解、讨论和演示的互动学习环境。",
        "bullets": [
            "核心卖点: 一键进入多智能体教学场景",
            "场景: 教学演示、课程陪练、知识讲解",
            "热度原因: 教育与 agent 交互深度结合",
        ],
        "narration": "OpenMAIC 展示了多智能体在教育场景的产品化方向，不只是聊天，而是完整课堂。",
        "duration": 10,
    },
    {
        "rank": 4,
        "repo": "karpathy/autoresearch",
        "language": "Python",
        "stars_today": 258,
        "stars_total": 42862,
        "headline": "让 AI 代理自己跑实验并迭代模型训练",
        "summary": "项目把单卡训练、实验记录和自动试错串起来，目标是隔夜产出更好的训练结果。",
        "bullets": [
            "核心卖点: agent 自动修改代码并评估实验",
            "场景: 小规模研究自动化、训练流程探索",
            "热度原因: Karpathy 亲自示范 AI 研究自动化",
        ],
        "narration": "Autoresearch 的亮点是把研究流程闭环化，让代理自己试验、保留结果并继续优化。",
        "duration": 10,
    },
    {
        "rank": 5,
        "repo": "msitarzewski/agency-agents",
        "language": "Shell",
        "stars_today": 259,
        "stars_total": 54829,
        "headline": "一个仓库打包整支 AI 专家团队",
        "summary": "用大量具有鲜明角色设定的专家代理覆盖设计、营销、开发与策略任务。",
        "bullets": [
            "核心卖点: 角色化专家 agent 集合",
            "场景: 内容、产品、增长、前端等协作",
            "热度原因: 人设鲜明且可直接复用",
        ],
        "narration": "Agency Agents 的传播力来自角色化设计，它让多人协作式 AI 工作流更容易想象。",
        "duration": 10,
    },
    {
        "rank": 6,
        "repo": "666ghj/MiroFish",
        "language": "Python",
        "stars_today": 230,
        "stars_total": 35117,
        "headline": "群体智能引擎尝试把预测问题统一抽象",
        "summary": "以 swarm intelligence 为核心，面向多种预测任务，强调通用性和多智能体推演能力。",
        "bullets": [
            "核心卖点: 通用群体智能预测引擎",
            "场景: 预测、推演、知识图谱辅助决策",
            "热度原因: AI 预测叙事和工程实现结合",
        ],
        "narration": "MiroFish 把多智能体预测包装成统一引擎，适合展示复杂推演类自动化场景。",
        "duration": 10,
    },
    {
        "rank": 7,
        "repo": "shareAI-lab/learn-claude-code",
        "language": "TypeScript",
        "stars_today": 234,
        "stars_total": 33023,
        "headline": "从零拆解 Claude Code 式 agent harness",
        "summary": "以教学方式复刻轻量 agent harness，帮助开发者理解命令行智能体的底层结构。",
        "bullets": [
            "核心卖点: 教学导向的 agent harness 拆解",
            "场景: 学习 CLI agent 架构与实现",
            "热度原因: 兼顾教育价值和实操价值",
        ],
        "narration": "Learn Claude Code 更像公开课，帮助团队理解如何自己搭出一套轻量智能体。",
        "duration": 10,
    },
    {
        "rank": 8,
        "repo": "ZhuLinsen/daily_stock_analysis",
        "language": "Python",
        "stars_today": 122,
        "stars_total": 23126,
        "headline": "把股票分析做成零成本定时运行的智能系统",
        "summary": "汇总行情、实时新闻和 LLM 决策面板，自动生成多市场分析与推送结果。",
        "bullets": [
            "核心卖点: 多数据源行情加 LLM 决策面板",
            "场景: A 股、港股、美股定时分析",
            "热度原因: 很接近真实可落地的个人自动化",
        ],
        "narration": "这个项目很适合自动化验证，因为它已经把数据抓取、分析和推送串成流程。",
        "duration": 10,
    },
    {
        "rank": 9,
        "repo": "gsd-build/get-shit-done",
        "language": "JavaScript",
        "stars_today": 230,
        "stars_total": 35293,
        "headline": "用上下文工程修复长会话里的性能衰减",
        "summary": "围绕 spec-driven development 和 context engineering，专门解决 agent 的上下文腐化问题。",
        "bullets": [
            "核心卖点: 元提示与上下文治理系统",
            "场景: 长任务编排、规范驱动开发",
            "热度原因: 命中所有 agent 用户的痛点",
        ],
        "narration": "Get Shit Done 关注的不是单次提示，而是长流程里如何持续保持上下文质量。",
        "duration": 10,
    },
    {
        "rank": 10,
        "repo": "opendataloader-project/opendataloader-pdf",
        "language": "Java",
        "stars_today": 199,
        "stars_total": 4672,
        "headline": "把 PDF 解析和可访问性自动化做到 AI 就绪",
        "summary": "支持结构化提取、版面分析、自动标签化和多语言 SDK，是文档智能处理的重要基础层。",
        "bullets": [
            "核心卖点: 面向 AI 的 PDF 结构化解析",
            "场景: RAG、文档抽取、可访问性合规",
            "热度原因: 企业文档场景需求非常明确",
        ],
        "narration": "OpenDataLoader PDF 是自动化流程里的基础设施，适合接在文档理解和知识抽取前面。",
        "duration": 10,
    },
]


def apply_gradient(slide):
    bg = slide.background.fill
    bg.gradient()
    stops = bg.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = BG_TOP
    stops[1].position = 1.0
    stops[1].color.rgb = BG_BOTTOM
    bg.gradient_angle = 90


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    size,
    color,
    bold=False,
    font_name="Microsoft YaHei",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.08,
):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_card(slide, left, top, width, height, fill_rgb, transparency=0.0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.fill.transparency = transparency
    shape.line.color.rgb = fill_rgb
    return shape


def add_chip(slide, left, top, width, text, *, fill_rgb, text_rgb=WHITE):
    chip = add_card(slide, left, top, width, Inches(0.42), fill_rgb, transparency=0.0)
    chip.line.color.rgb = fill_rgb
    add_textbox(
        slide,
        left + Inches(0.03),
        top + Inches(0.04),
        width - Inches(0.06),
        Inches(0.22),
        text,
        size=11,
        color=text_rgb,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    return chip


def parse_feature_cards(item: dict) -> list[dict]:
    cards = [{"icon": "🧩", "title": "项目简介", "body": item["summary"]}]
    for bullet in item["bullets"]:
        if ":" in bullet:
            label, body = bullet.split(":", 1)
            label = label.strip()
            body = body.strip()
        else:
            label = "项目定位"
            body = bullet.strip()
        icon, title = CARD_META.get(label, ("✨", label))
        cards.append({"icon": icon, "title": title, "body": body})
    return cards[:4]


def format_timecode(total_seconds: int) -> str:
    minutes, seconds = divmod(total_seconds, 60)
    return f"{minutes:02d}:{seconds:02d}"


def write_tts_timeline(markdown_path: Path, slides: list[dict]) -> None:
    lines = [
        "# TTS Timeline",
        "",
        "用于后续 TTS 调试与音视频对齐。",
        "",
    ]
    cursor = 0
    for slide in slides:
        start = cursor
        end = cursor + slide["duration"]
        lines.extend(
            [
                f"## Slide {slide['index']:02d} - {slide['title']}",
                "",
                f"- Start: `{format_timecode(start)}`",
                f"- End: `{format_timecode(end)}`",
                f"- Duration: `{slide['duration']}s`",
                f"- Suggested audio: `audio/slide-{slide['index']:02d}.mp3`",
                f"- Voiceover: {slide['narration']}",
                "",
            ]
        )
        cursor = end
    markdown_path.write_text("\n".join(lines), encoding="utf-8")


def build_title_slide(prs: PresentationType):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_gradient(slide)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(8.1), Inches(0.6), Inches(4.2), Inches(4.2)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = CYAN
    accent.fill.transparency = 0.74
    accent.line.color.rgb = CYAN

    accent2 = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(-0.4), Inches(4.2), Inches(3.1), Inches(3.1)
    )
    accent2.fill.solid()
    accent2.fill.fore_color.rgb = GOLD
    accent2.fill.transparency = 0.78
    accent2.line.color.rgb = GOLD

    tag = add_card(slide, Inches(0.7), Inches(0.65), Inches(2.3), Inches(0.5), BLUE)
    tag.fill.transparency = 0.08
    add_textbox(
        slide,
        Inches(0.82),
        Inches(0.73),
        Inches(2.0),
        Inches(0.3),
        "AUTOMATED VIDEO PPTX",
        size=14,
        color=WHITE,
        bold=True,
    )

    add_textbox(
        slide,
        Inches(0.7),
        Inches(1.35),
        Inches(7.5),
        Inches(1.4),
        "今日最火 10 个开源项目",
        size=30,
        color=WHITE,
        bold=True,
    )
    add_textbox(
        slide,
        Inches(0.7),
        Inches(2.28),
        Inches(7.0),
        Inches(0.9),
        "卡片式结构化速览 · 适配自动播片、TTS 时间轴与后续视频合成",
        size=17,
        color=WHITE,
    )

    card_specs = [
        (
            Inches(0.75),
            "01",
            "采集来源",
            "OSS Insight 过去 24 小时趋势榜，聚焦当天最热开源项目。",
        ),
        (
            Inches(4.45),
            "02",
            "页面结构",
            "每页包含编号、项目名、项目简介和多张特性卡片，不显示口播稿。",
        ),
        (
            Inches(8.15),
            "03",
            "输出产物",
            "同步生成 PPTX、MP4、JSON manifest 和 Markdown 时间轴文件。",
        ),
    ]
    for left, step, title, body in card_specs:
        card = add_card(slide, left, Inches(3.65), Inches(3.45), Inches(2.08), WHITE)
        card.fill.transparency = 0.08
        card.line.color.rgb = WHITE
        add_chip(
            slide, left + Inches(0.18), Inches(3.86), Inches(0.72), step, fill_rgb=BLUE
        )
        add_textbox(
            slide,
            left + Inches(0.18),
            Inches(4.34),
            Inches(2.95),
            Inches(0.35),
            title,
            size=18,
            color=NAVY,
            bold=True,
        )
        add_textbox(
            slide,
            left + Inches(0.18),
            Inches(4.8),
            Inches(2.95),
            Inches(0.7),
            body,
            size=13,
            color=TEXT,
        )

    add_textbox(
        slide,
        Inches(0.78),
        Inches(6.45),
        Inches(3.6),
        Inches(0.35),
        f"生成日期: {date.today().isoformat()}",
        size=12,
        color=WHITE,
    )


def build_project_slide(prs: PresentationType, item: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    apply_gradient(slide)
    cards = parse_feature_cards(item)

    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.34),
    ).fill.solid()
    topbar = slide.shapes[0]
    topbar.fill.fore_color.rgb = BLUE
    topbar.line.color.rgb = BLUE

    rank_card = add_card(
        slide, Inches(0.72), Inches(0.62), Inches(1.25), Inches(1.1), NAVY
    )
    rank_card.fill.transparency = 0.0
    add_textbox(
        slide,
        Inches(0.95),
        Inches(0.82),
        Inches(0.7),
        Inches(0.45),
        f"{item['rank']:02d}",
        size=24,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(0.78),
        Inches(1.28),
        Inches(1.12),
        Inches(0.2),
        "TOP ITEM",
        size=8,
        color=WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        Inches(2.15),
        Inches(0.62),
        Inches(6.5),
        Inches(0.6),
        item["repo"],
        size=24,
        color=WHITE,
        bold=True,
    )
    add_textbox(
        slide,
        Inches(2.15),
        Inches(1.08),
        Inches(7.2),
        Inches(0.52),
        item["headline"],
        size=15,
        color=WHITE,
    )
    add_textbox(
        slide,
        Inches(2.15),
        Inches(1.48),
        Inches(6.7),
        Inches(0.42),
        f"项目介绍: {item['summary']}",
        size=13,
        color=WHITE,
    )

    add_chip(
        slide, Inches(9.05), Inches(0.72), Inches(1.15), item["language"], fill_rgb=CYAN
    )
    add_chip(
        slide,
        Inches(10.35),
        Inches(0.72),
        Inches(1.05),
        f"+{item['stars_today']}",
        fill_rgb=GOLD,
        text_rgb=NAVY,
    )
    add_chip(
        slide,
        Inches(11.52),
        Inches(0.72),
        Inches(1.0),
        f"{item['stars_total']}",
        fill_rgb=WHITE,
        text_rgb=NAVY,
    )

    grid_positions = [
        (Inches(0.78), Inches(2.2)),
        (Inches(6.75), Inches(2.2)),
        (Inches(0.78), Inches(4.28)),
        (Inches(6.75), Inches(4.28)),
    ]
    for card_item, (left, top) in zip(cards, grid_positions):
        card = add_card(slide, left, top, Inches(5.78), Inches(1.74), WHITE)
        card.fill.transparency = 0.03
        card.line.color.rgb = SOFT
        add_textbox(
            slide,
            left + Inches(0.22),
            top + Inches(0.18),
            Inches(0.42),
            Inches(0.28),
            card_item["icon"],
            size=18,
            color=BLUE,
            bold=True,
            font_name=ICON_FONT,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left + Inches(0.72),
            top + Inches(0.14),
            Inches(3.5),
            Inches(0.3),
            card_item["title"],
            size=15,
            color=NAVY,
            bold=True,
        )
        add_textbox(
            slide,
            left + Inches(0.22),
            top + Inches(0.6),
            Inches(5.2),
            Inches(0.82),
            card_item["body"],
            size=15,
            color=TEXT,
        )

    footer = add_card(
        slide, Inches(0.78), Inches(6.28), Inches(11.75), Inches(0.52), WHITE
    )
    footer.fill.transparency = 0.05
    footer.line.color.rgb = SOFT
    add_textbox(
        slide,
        Inches(1.0),
        Inches(6.4),
        Inches(1.2),
        Inches(0.2),
        "项目结论",
        size=12,
        color=BLUE,
        bold=True,
    )
    add_textbox(
        slide,
        Inches(2.15),
        Inches(6.37),
        Inches(9.95),
        Inches(0.24),
        item["headline"],
        size=14,
        color=TEXT,
        bold=True,
    )


def export_video_via_powershell(
    pptx_path: Path, video_path: Path, durations: list[int]
):
    durations_text = ",".join(str(value) for value in durations)
    script = f"""
$ErrorActionPreference = 'Stop'
$pptPath = '{pptx_path.resolve()}'
$videoPath = '{video_path.resolve()}'
$durations = @({durations_text})
    $app = New-Object -ComObject PowerPoint.Application
    $presentation = $app.Presentations.Open($pptPath, $false, $false, $false)
for ($i = 1; $i -le $presentation.Slides.Count; $i++) {{
  $slide = $presentation.Slides.Item($i)
  $slide.SlideShowTransition.AdvanceOnTime = $true
  $slide.SlideShowTransition.AdvanceTime = [double]$durations[$i - 1]
}}
$presentation.Save()
$presentation.CreateVideo($videoPath, $true, 10, 1080, 30, 85)
$deadline = (Get-Date).AddMinutes(20)
while ((Get-Date) -lt $deadline) {{
  Start-Sleep -Seconds 2
  $status = $presentation.CreateVideoStatus
  if ($status -eq 3) {{
    break
  }}
  if ($status -eq 4) {{
    throw 'PowerPoint video export failed.'
  }}
}}
if ($presentation.CreateVideoStatus -ne 3) {{
  throw 'PowerPoint video export timed out.'
}}
$presentation.Close()
$app.Quit()
""".strip()
    encoded = base64.b64encode(script.encode("utf-16le")).decode("ascii")
    subprocess.run(
        [
            "powershell",
            "-NoProfile",
            "-EncodedCommand",
            encoded,
        ],
        check=True,
    )


def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_title_slide(prs)
    for item in SLIDES:
        build_project_slide(prs, item)

    prs.save(str(PPTX_PATH))

    manifest = {
        "generated_on": date.today().isoformat(),
        "source": "https://api.ossinsight.io/v1/trends/repos/?period=past_24_hours&language=All",
        "video_seconds": 6 + sum(item["duration"] for item in SLIDES),
        "tts_timeline_markdown": str(TIMELINE_MD_PATH),
        "slides": [
            {
                "index": 1,
                "title": "今日最火 10 个开源项目",
                "duration": 6,
                "narration": "这是一个用于验证 PPTX 自动转视频的样片，内容来自今天的开源趋势榜。",
            },
            *[
                {
                    "index": idx + 2,
                    "title": item["repo"],
                    "duration": item["duration"],
                    "narration": item["narration"],
                }
                for idx, item in enumerate(SLIDES)
            ],
        ],
    }
    MANIFEST_PATH.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8"
    )
    write_tts_timeline(TIMELINE_MD_PATH, manifest["slides"])

    export_video_via_powershell(
        PPTX_PATH, VIDEO_PATH, [6] + [item["duration"] for item in SLIDES]
    )
    print(f"Generated {PPTX_PATH}")
    print(f"Generated {VIDEO_PATH}")
    print(f"Generated {MANIFEST_PATH}")
    print(f"Generated {TIMELINE_MD_PATH}")


if __name__ == "__main__":
    main()
